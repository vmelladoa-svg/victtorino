"""Genera presentacion_victoria_2026-05-18.pptx (12 slides)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = r"C:\Users\dell\victtorino\presentacion_victoria_2026-05-18.pptx"

# Paleta
NAVY     = RGBColor(0x0F, 0x17, 0x2A)
BLUE     = RGBColor(0x1E, 0x40, 0xAF)
BLUE_LT  = RGBColor(0x3B, 0x82, 0xF6)
ACCENT   = RGBColor(0xF9, 0x73, 0x16)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
GRAY     = RGBColor(0x64, 0x74, 0x8B)
GRAY_LT  = RGBColor(0xE2, 0xE8, 0xF0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BG_SOFT  = RGBColor(0xF8, 0xFA, 0xFC)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=18, color=NAVY, bullet_color=BLUE_LT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        r1 = p.add_run()
        r1.text = "•  "
        r1.font.size = Pt(size + 2)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r1.font.name = font
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = font
    return tb


def add_accent_bar(slide, x, y, w=Inches(0.15), h=Inches(0.6), color=BLUE_LT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def add_card(slide, x, y, w, h, fill=BG_SOFT, line=GRAY_LT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.06
    card.fill.solid(); card.fill.fore_color.rgb = fill
    card.line.color.rgb = line
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    return card


def slide_title_block(slide, title, kicker=None):
    add_accent_bar(slide, Inches(0.6), Inches(0.55), h=Inches(0.55))
    add_text(slide, Inches(0.85), Inches(0.45), Inches(11.5), Inches(0.8),
             title, size=32, bold=True, color=NAVY)
    if kicker:
        add_text(slide, Inches(0.85), Inches(1.05), Inches(11.5), Inches(0.4),
                 kicker, size=14, color=GRAY)


# ───────── Slide 1: Portada ─────────
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
# acento decorativo
deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.4), Inches(0.4), Inches(0.7))
deco.fill.solid(); deco.fill.fore_color.rgb = BLUE_LT
deco.line.fill.background(); deco.shadow.inherit = False
add_text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(0.6),
         "VICTTORINO  ·  Mayo 2026", size=14, color=BLUE_LT, font="Calibri", bold=True)
add_text(s, Inches(0.6), Inches(3.2), Inches(12), Inches(1.4),
         "Victoria", size=72, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.8),
         "Asistente Virtual y Sistema de Automatización", size=28, color=GRAY_LT)
add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
         "Estado del proyecto · Equipo interno", size=14, color=GRAY)

# ───────── Slide 2: El desafío ─────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title_block(s, "El desafío", "Por qué necesitamos automatizar")
add_bullets(s, Inches(0.85), Inches(1.8), Inches(11.5), Inches(5),
            [
                "Operamos en 6 canales: tienda física, victtorino.cl, MercadoLibre (3 cuentas), Falabella, París y Walmart.",
                "Decenas de preguntas y pedidos por día — cada uno necesita respuesta clara y rápida.",
                "Atender manualmente no escala: noche, fines de semana, vacaciones.",
                "Cada minuto sin respuesta es una venta que se enfría o que se va a la competencia.",
                "Necesitamos respuestas consistentes, en chileno, con la información correcta — siempre.",
            ], size=20)

# ───────── Slide 3: La solución ─────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title_block(s, "Lo que construimos", "Tres piezas que trabajan juntas")
# 3 cards
card_w = Inches(3.9); card_h = Inches(4); top = Inches(2.2)
labels = [
    ("01", "Victoria",
     "Asistente de WhatsApp con IA. Habla chileno, conoce el negocio y responde como un buen vendedor.", BLUE_LT),
    ("02", "Automatización ML",
     "Las 3 cuentas de MercadoLibre responden preguntas y envían mensajes post-compra automáticamente.", ACCENT),
    ("03", "Operación de fondo",
     "Auditorías, mejora de fichas, reposición de stock, informes — todo en piloto automático.", GREEN),
]
for i, (num, t, desc, color) in enumerate(labels):
    x = Inches(0.6 + i * 4.25)
    add_card(s, x, top, card_w, card_h)
    add_accent_bar(s, x + Inches(0.4), top + Inches(0.4), w=Inches(0.35), h=Inches(0.35), color=color)
    add_text(s, x + Inches(0.4), top + Inches(0.85), card_w - Inches(0.5), Inches(0.5),
             num, size=14, bold=True, color=color)
    add_text(s, x + Inches(0.4), top + Inches(1.25), card_w - Inches(0.5), Inches(0.7),
             t, size=24, bold=True, color=NAVY)
    add_text(s, x + Inches(0.4), top + Inches(2.05), card_w - Inches(0.7), Inches(1.8),
             desc, size=15, color=GRAY)

# ───────── Slide 4: Victoria WhatsApp ─────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title_block(s, "Victoria — Atención por WhatsApp", "El agente que conversa por nosotros 24/7")
add_bullets(s, Inches(0.85), Inches(1.85), Inches(7.5), Inches(5),
            [
                "Habla español chileno, tono amigable con cercanía profesional.",
                "Sabe productos, precios top, envíos a regiones y medios de pago.",
                "Atiende dudas, califica ventas y orienta a victtorino.cl o a la tienda física.",
                "Detecta cuándo un humano respondió: se queda callada por 60 min para no pisar.",
                "Cada conversación queda registrada en memoria — el cliente no se repite.",
            ], size=18)
# panel lateral derecho — ejemplo de conversación
panel_x = Inches(8.7); panel_y = Inches(1.85); panel_w = Inches(4); panel_h = Inches(5.2)
add_card(s, panel_x, panel_y, panel_w, panel_h, fill=BG_SOFT)
add_text(s, panel_x + Inches(0.25), panel_y + Inches(0.2), panel_w - Inches(0.4), Inches(0.4),
         "Ejemplo de conversación", size=12, bold=True, color=GRAY)
# burbuja cliente
b1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, panel_x + Inches(0.25), panel_y + Inches(0.7), Inches(3), Inches(0.7))
b1.adjustments[0] = 0.3
b1.fill.solid(); b1.fill.fore_color.rgb = GRAY_LT
b1.line.fill.background(); b1.shadow.inherit = False
add_text(s, panel_x + Inches(0.4), panel_y + Inches(0.78), Inches(2.7), Inches(0.6),
         "Tienen lavaplatos 80x44 con secador?", size=12, color=NAVY)
# burbuja Victoria
b2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, panel_x + Inches(0.75), panel_y + Inches(1.55), Inches(3), Inches(1.8))
b2.adjustments[0] = 0.18
b2.fill.solid(); b2.fill.fore_color.rgb = BLUE_LT
b2.line.fill.background(); b2.shadow.inherit = False
add_text(s, panel_x + Inches(0.9), panel_y + Inches(1.65), Inches(2.7), Inches(1.7),
         "¡Hola! Sí, tenemos el de 80x44 con secador derecho o izquierdo a $33.082. "
         "Lo puedes ver en victtorino.cl. ¿Necesitas envío?",
         size=12, color=WHITE)
# burbuja cliente 2
b3 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, panel_x + Inches(0.25), panel_y + Inches(3.55), Inches(3), Inches(0.7))
b3.adjustments[0] = 0.3
b3.fill.solid(); b3.fill.fore_color.rgb = GRAY_LT
b3.line.fill.background(); b3.shadow.inherit = False
add_text(s, panel_x + Inches(0.4), panel_y + Inches(3.62), Inches(2.7), Inches(0.6),
         "Sí, a Concepción", size=12, color=NAVY)
# burbuja Victoria 2
b4 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, panel_x + Inches(0.75), panel_y + Inches(4.4), Inches(3), Inches(0.65))
b4.adjustments[0] = 0.3
b4.fill.solid(); b4.fill.fore_color.rgb = BLUE_LT
b4.line.fill.background(); b4.shadow.inherit = False
add_text(s, panel_x + Inches(0.9), panel_y + Inches(4.48), Inches(2.7), Inches(0.6),
         "Envío $4.990, gratis sobre $60.000 😊", size=12, color=WHITE)

# ───────── Slide 5: MercadoLibre ─────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title_block(s, "MercadoLibre — 3 cuentas en piloto automático",
                  "Preguntas + post-compra, sin levantar un dedo")
# tabla simple con las 3 cuentas
cuentas = [
    ("Cuenta", "Alias", "Lo que hace"),
    ("C1", "PREMIUMGRIFERIAS1", "Responde preguntas + mensaje post-compra"),
    ("C2", "VICTTORINOFICIAL2", "Responde preguntas + mensaje post-compra"),
    ("C3", "NOVAGRIFERIAS3",    "Responde preguntas + mensaje post-compra"),
]
tbl_x = Inches(0.85); tbl_y = Inches(1.95)
rows = len(cuentas); cols = 3
col_widths = [Inches(1.2), Inches(3.5), Inches(7.8)]
row_h = Inches(0.55)
table_shape = s.shapes.add_table(rows, cols, tbl_x, tbl_y, sum(col_widths, Emu(0)), row_h * rows).table
for ci, cw in enumerate(col_widths):
    table_shape.columns[ci].width = cw
for ri, row in enumerate(cuentas):
    for ci, val in enumerate(row):
        cell = table_shape.cell(ri, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = val
        r.font.size = Pt(16); r.font.name = "Calibri"
        r.font.bold = (ri == 0)
        r.font.color.rgb = WHITE if ri == 0 else NAVY
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if ri == 0 else (BG_SOFT if ri % 2 else WHITE)
# resumen abajo
add_bullets(s, Inches(0.85), Inches(4.7), Inches(11.5), Inches(2.5),
            [
                "Polling cada 5 min: detecta preguntas nuevas y responde con la información del producto.",
                "Polling cada 10 min: detecta órdenes nuevas (últimas 24h) y envía mensaje post-compra al comprador.",
                "Mensaje post-compra sugiere productos complementarios según el producto comprado.",
                "Cada acción se notifica al dueño por WhatsApp.",
            ], size=16)

# ───────── Slide 6: Operación de fondo ─────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title_block(s, "Lo que sucede sin que te enteres", "Tareas automáticas y herramientas operativas")
items_op = [
    ("Auditorías", "Calidad de fichas en las 3 cuentas ML. Detecta fichas con atributos incompletos, fotos malas, descripciones flojas."),
    ("Mejora de fichas", "Genera descripciones e imágenes mejoradas con IA y las publica."),
    ("Reposición de stock", "Mantiene el inventario sincronizado entre canales."),
    ("Sync WooCommerce", "victtorino.cl y el sistema central comparten productos y pedidos."),
    ("Informes ejecutivos", "PDFs gerenciales generados automáticamente con métricas del negocio."),
    ("Tareas nocturnas", "Refresh de tokens, scrapers de facturación, reportes — todo programado."),
]
top = Inches(1.9); card_w = Inches(6); card_h = Inches(1.55)
for i, (t, desc) in enumerate(items_op):
    col = i % 2; row = i // 2
    x = Inches(0.6 + col * 6.3); y = top + Inches(row * 1.7)
    add_card(s, x, y, card_w, card_h)
    add_accent_bar(s, x + Inches(0.25), y + Inches(0.3), w=Inches(0.1), h=Inches(0.5), color=BLUE_LT)
    add_text(s, x + Inches(0.5), y + Inches(0.22), card_w - Inches(0.7), Inches(0.45),
             t, size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(0.5), y + Inches(0.7), card_w - Inches(0.7), Inches(0.85),
             desc, size=13, color=GRAY)

# ───────── Slide 7: Arquitectura ─────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title_block(s, "Cómo se conecta todo", "Una vista simple del sistema")
# fila clientes
clientes = [("WhatsApp\n(clientes)", BLUE_LT), ("MercadoLibre\n(3 cuentas)", ACCENT), ("victtorino.cl\n(WooCommerce)", GREEN)]
for i, (label, c) in enumerate(clientes):
    x = Inches(1.2 + i * 4); y = Inches(1.95)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3), Inches(1))
    box.adjustments[0] = 0.18
    box.fill.solid(); box.fill.fore_color.rgb = c
    box.line.fill.background(); box.shadow.inherit = False
    add_text(s, x, y, Inches(3), Inches(1), label, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
# flechas hacia el centro
for i in range(3):
    x = Inches(2.7 + i * 4); y1 = Inches(2.95); y2 = Inches(3.5)
    ln = s.shapes.add_connector(1, x, y1, x, y2)
    ln.line.color.rgb = GRAY
    ln.line.width = Pt(2)
# Victoria al centro
victoria = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(3.5), Inches(5.3), Inches(1.5))
victoria.adjustments[0] = 0.15
victoria.fill.solid(); victoria.fill.fore_color.rgb = NAVY
victoria.line.fill.background(); victoria.shadow.inherit = False
add_text(s, Inches(4), Inches(3.55), Inches(5.3), Inches(0.5),
         "VICTORIA", size=14, bold=True, color=BLUE_LT, align=PP_ALIGN.CENTER)
add_text(s, Inches(4), Inches(3.95), Inches(5.3), Inches(0.5),
         "Cerebro IA (Claude)", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(4), Inches(4.45), Inches(5.3), Inches(0.5),
         "FastAPI · Memoria de conversaciones", size=12, color=GRAY_LT, align=PP_ALIGN.CENTER)
# flecha hacia DB
ln = s.shapes.add_connector(1, Inches(6.65), Inches(5), Inches(6.65), Inches(5.5))
ln.line.color.rgb = GRAY; ln.line.width = Pt(2)
# DB
db = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(5.5), Inches(5.3), Inches(1))
db.adjustments[0] = 0.18
db.fill.solid(); db.fill.fore_color.rgb = BG_SOFT
db.line.color.rgb = GRAY_LT; db.shadow.inherit = False
add_text(s, Inches(4), Inches(5.5), Inches(5.3), Inches(1),
         "Base de datos · Railway", size=16, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
# nota inferior
add_text(s, Inches(0.85), Inches(6.85), Inches(11.5), Inches(0.5),
         "Todo corre en la nube (Railway). Sin servidores propios, sin mantención.",
         size=12, color=GRAY)

# ───────── Slide 8: Estado actual ─────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title_block(s, "Hoy en producción", "Lo que está vivo y funcionando")
data = [
    ("Componente", "Estado"),
    ("Victoria en Railway", "Activo 24/7"),
    ("Atención WhatsApp (Meta API)", "Activa"),
    ("Preguntas MercadoLibre (3 cuentas)", "Activa, polling 5 min"),
    ("Mensajes post-compra", "Activa, polling 10 min"),
    ("Renovación de tokens ML", "Automática cada 5 horas"),
    ("Sincronización con WooCommerce", "Activa"),
    ("Logs y memoria persistente", "Activa"),
]
rows = len(data); cols = 2
tx = Inches(0.85); ty = Inches(1.9)
cw = [Inches(7), Inches(4.5)]
rh = Inches(0.55)
t = s.shapes.add_table(rows, cols, tx, ty, sum(cw, Emu(0)), rh * rows).table
for ci, w in enumerate(cw):
    t.columns[ci].width = w
for ri, row in enumerate(data):
    for ci, val in enumerate(row):
        cell = t.cell(ri, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        # icono check si es columna estado y no header
        if ci == 1 and ri > 0:
            r0 = p.add_run(); r0.text = "✓  "
            r0.font.size = Pt(16); r0.font.bold = True
            r0.font.color.rgb = GREEN; r0.font.name = "Calibri"
        r = p.add_run(); r.text = val
        r.font.size = Pt(16); r.font.name = "Calibri"
        r.font.bold = (ri == 0)
        r.font.color.rgb = WHITE if ri == 0 else NAVY
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if ri == 0 else (BG_SOFT if ri % 2 else WHITE)

# ───────── Slide 9: Beneficios ─────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title_block(s, "Lo que estamos ganando", "Valor que se siente en el día a día")
benefits = [
    ("⏱", "Tiempo", "Respuestas instantáneas. Sin esperar al horario laboral. Sin perder ventas en la noche o fines de semana."),
    ("⚡", "Escala", "Atiende N consultas en paralelo sin colas. Si crece el volumen, no hay que contratar más gente."),
    ("✓", "Consistencia", "Mismo tono, misma información, mismas políticas — siempre. Cero errores por cansancio o distracción."),
    ("📊", "Datos", "Cada conversación queda registrada. Patrones, dudas frecuentes, productos que más se preguntan — todo medible."),
]
top = Inches(1.95); cw = Inches(6); ch = Inches(2.4)
for i, (icon, t, desc) in enumerate(benefits):
    col = i % 2; row = i // 2
    x = Inches(0.6 + col * 6.3); y = top + Inches(row * 2.55)
    add_card(s, x, y, cw, ch)
    add_text(s, x + Inches(0.35), y + Inches(0.3), Inches(1), Inches(1),
             icon, size=32, color=BLUE_LT)
    add_text(s, x + Inches(1.3), y + Inches(0.35), cw - Inches(1.5), Inches(0.55),
             t, size=22, bold=True, color=NAVY)
    add_text(s, x + Inches(1.3), y + Inches(0.95), cw - Inches(1.5), Inches(1.4),
             desc, size=13, color=GRAY)

# ───────── Slide 10: Roadmap ─────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title_block(s, "Lo que falta", "Próximas mejoras priorizadas")
roadmap = [
    ("Corto plazo (1-2 meses)",
     ["Dashboard de métricas en vivo (cuántas conversaciones, ventas atribuidas, tiempo de respuesta).",
      "Que Victoria consulte stock real en tiempo real antes de responder.",
      "Tests A/B de mensajes post-compra para optimizar la conversión."]),
    ("Mediano plazo (3-6 meses)",
     ["Integrar Falabella, París y Walmart con la misma lógica.",
      "Que Victoria pueda crear pedidos en victtorino.cl directamente desde WhatsApp.",
      "Audios de WhatsApp (hoy solo procesa texto)."]),
    ("Largo plazo (6+ meses)",
     ["Predicción de demanda y sugerencia de reposición.",
      "Asistente de operaciones — Victor consulta el estado del negocio por WhatsApp.",
      "Atribución de ventas: saber qué conversación cerró cada venta."]),
]
y = Inches(1.85)
for label, items in roadmap:
    add_text(s, Inches(0.85), y, Inches(11.5), Inches(0.4),
             label, size=15, bold=True, color=BLUE_LT)
    add_bullets(s, Inches(1.1), y + Inches(0.45), Inches(11), Inches(1.3),
                items, size=13)
    y += Inches(1.7)

# ───────── Slide 11: Visión ─────────
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_accent_bar(s, Inches(0.6), Inches(1.5), h=Inches(0.55), color=BLUE_LT)
add_text(s, Inches(0.85), Inches(1.4), Inches(11.5), Inches(0.7),
         "Hacia dónde vamos", size=18, color=BLUE_LT, bold=True)
add_text(s, Inches(0.85), Inches(2.5), Inches(11.7), Inches(3),
         "Un agente comercial que opere los canales\ndigitales como nuestro mejor vendedor —",
         size=42, bold=True, color=WHITE)
add_text(s, Inches(0.85), Inches(4.5), Inches(11.7), Inches(1.5),
         "24/7, en todos los canales, con datos para decidir.",
         size=32, color=BLUE_LT, bold=True)
add_text(s, Inches(0.85), Inches(6.4), Inches(11.7), Inches(0.6),
         "Liberar tiempo del equipo. Crecer sin que crezca la carga operativa.",
         size=16, color=GRAY_LT)

# ───────── Slide 12: Cierre ─────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_accent_bar(s, Inches(0.6), Inches(2.7), h=Inches(0.7), color=BLUE_LT)
add_text(s, Inches(0.85), Inches(2.55), Inches(11.7), Inches(1.2),
         "¿Preguntas?", size=60, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(3.9), Inches(11.7), Inches(0.7),
         "Conversemos.", size=24, color=GRAY)
add_text(s, Inches(0.85), Inches(6.5), Inches(11.7), Inches(0.5),
         "Victoria · Victtorino · Mayo 2026", size=12, color=GRAY)

prs.save(OUT)
print("OK:", OUT)
